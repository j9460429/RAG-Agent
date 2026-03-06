"""Marker 文件解析微服務 — 多引擎文件解析。

- PDF/DOCX: 使用 pymupdf4llm（結構化 Markdown）
- PPTX: 使用 python-pptx（投影片結構 → Markdown 標題層級）
- XLSX/XLS: 使用 openpyxl（Markdown 表格格式）
- 其他格式: pymupdf4llm fallback
"""

import base64
import logging
import os
import sys
import tempfile
import time
import traceback
from typing import Any

from fastapi import FastAPI, File, UploadFile
from fastapi.responses import JSONResponse

# 設定日誌
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    stream=sys.stdout,
)
logger = logging.getLogger("marker-service")

app = FastAPI(title="Marker Document Parser", version="2.1.0")


@app.get("/health")
async def health():
    """健康檢查端點。"""
    return {"status": "ok", "service": "marker", "engine": "pymupdf4llm+pptx+xlsx"}


@app.post("/parse")
async def parse_document(file: UploadFile = File(...)):
    """
    解析上傳的文件，回傳 Markdown + 結構化 chunks。

    支援格式：PDF, DOCX, PPTX, XLSX, HTML, EPUB
    回傳：{ markdown, chunks[], metadata }
    """
    start_time = time.time()

    try:
        file_bytes = await file.read()
        filename = file.filename or "document.pdf"
        logger.info(f"收到解析請求: {filename} ({len(file_bytes)} bytes)")

        suffix = os.path.splitext(filename)[1].lower() or ".pdf"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name

        page_images = None
        try:
            if suffix == ".pptx":
                markdown_text, page_count = _parse_pptx(tmp_path)
            elif suffix in (".xlsx", ".xls"):
                markdown_text, page_count = _parse_xlsx(tmp_path)
            else:
                markdown_text, page_count, page_images = _parse_with_pymupdf(tmp_path)
        finally:
            os.unlink(tmp_path)

        logger.info(f"解析完成: {len(markdown_text)} chars, {page_count} pages, images={'yes' if page_images else 'no'}")

        chunks = _build_chunks(markdown_text)

        elapsed = time.time() - start_time

        response_content: dict[str, Any] = {
            "success": True,
            "markdown": markdown_text,
            "chunks": chunks,
            "metadata": {
                "filename": filename,
                "page_count": page_count,
                "parse_time_seconds": round(elapsed, 2),
            },
        }

        # 圖片型 PDF：附帶頁面圖片供 OCR
        if page_images:
            response_content["page_images"] = page_images
            response_content["metadata"]["ocr_needed"] = True

        return JSONResponse(content=response_content)

    except Exception as e:
        traceback.print_exc()
        return JSONResponse(
            status_code=500,
            content={
                "success": False,
                "error": str(e),
                "error_type": type(e).__name__,
            },
        )


def _parse_xlsx(file_path: str) -> tuple[str, int]:
    """使用 openpyxl 解析 Excel，輸出 Markdown 表格格式。

    每個工作表 → 一個 `## 工作表: 名稱` 區塊。
    資料列轉為 Markdown 表格（| col1 | col2 | ...）。
    """
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    sheet_count = len(wb.sheetnames)
    parts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[list[str]] = []

        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            # 跳過全空行
            if not any(cells):
                continue
            rows.append(cells)

        if not rows:
            continue

        # 建構 Markdown 表格
        header = f"## 工作表: {sheet_name}\n\n"

        # 第一行為表頭
        col_count = max(len(r) for r in rows)
        # 補齊欄位數
        normalized = [r + [""] * (col_count - len(r)) for r in rows]

        table_header = "| " + " | ".join(normalized[0]) + " |"
        separator = "| " + " | ".join(["---"] * col_count) + " |"
        table_rows = [
            "| " + " | ".join(r) + " |" for r in normalized[1:]
        ]

        table_md = "\n".join([table_header, separator, *table_rows])
        parts.append(header + table_md)

    wb.close()

    if not parts:
        return "", 0

    markdown = "\n\n---\n\n".join(parts)
    return markdown, sheet_count


def _parse_pptx(file_path: str) -> tuple[str, int]:
    """使用 python-pptx 解析 PPTX，保留投影片結構。

    每張投影片 → 一個 `## 投影片 N: 標題` 區塊。
    投影片內的文字框按段落層級產出 Markdown。
    """
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation(file_path)
    page_count = len(prs.slides)
    parts: list[str] = []

    for idx, slide in enumerate(prs.slides, 1):
        slide_title = ""
        slide_body_parts: list[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            # 取得投影片標題
            if slide.shapes.title and shape.shape_id == slide.shapes.title.shape_id:
                slide_title = shape.text_frame.text.strip()
                continue

            # 處理文字框內容
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                # 根據段落層級產出 Markdown
                level = para.level or 0
                font_size = None
                if para.runs:
                    font_size = para.runs[0].font.size

                if font_size and font_size >= Pt(20):
                    slide_body_parts.append(f"### {text}")
                elif level == 0:
                    slide_body_parts.append(text)
                else:
                    indent = "  " * level
                    slide_body_parts.append(f"{indent}- {text}")

        # 組裝投影片區塊
        header = f"## 投影片 {idx}"
        if slide_title:
            header = f"## 投影片 {idx}: {slide_title}"

        slide_content = header + "\n\n"
        if slide_body_parts:
            slide_content += "\n\n".join(slide_body_parts) + "\n"

        parts.append(slide_content)

    markdown = "\n\n---\n\n".join(parts)
    return markdown, page_count


def _parse_with_pymupdf(file_path: str) -> tuple[str, int, list[str] | None]:
    """使用 pymupdf4llm 解析 PDF/DOCX 等格式。

    當文字提取為空時（圖片型 PDF），自動將每頁轉為 PNG base64 回傳，
    供 app 端送 Gemini Vision OCR。

    Returns:
        (markdown_text, page_count, page_images)
        page_images 為 None 表示有文字；為 list[str] 表示圖片型 PDF。
    """
    import pymupdf
    import pymupdf4llm

    doc = pymupdf.open(file_path)
    page_count = len(doc)
    doc.close()

    markdown_text = pymupdf4llm.to_markdown(
        file_path,
        show_progress=False,
    )

    # 文字夠多 → 正常回傳
    if markdown_text.strip():
        return markdown_text, page_count, None

    # 空文字 → 圖片型 PDF，把每頁轉為 base64 PNG
    logger.info(f"文字為空，轉為圖片模式（{page_count} 頁）")
    doc = pymupdf.open(file_path)
    page_images: list[str] = []
    for page in doc:
        # 100 DPI — OCR 足夠清晰，大幅縮小 payload（比 150 DPI 節省約 56% 大小）
        pix = page.get_pixmap(dpi=100)
        png_bytes = pix.tobytes("png")
        b64 = base64.b64encode(png_bytes).decode("ascii")
        page_images.append(b64)
    doc.close()

    return "", page_count, page_images


def _build_chunks(
    markdown: str,
) -> list[dict[str, Any]]:
    """
    將 Markdown 文本切分為結構化 chunks。

    策略：按 heading 分段，每段標記 chunk_type。
    """
    if not markdown or not markdown.strip():
        return []

    chunks: list[dict[str, Any]] = []
    current_text = ""
    current_type = "text"
    current_page = 1

    for line in markdown.split("\n"):
        stripped = line.strip()

        # 偵測 heading
        if stripped.startswith("#"):
            if current_text.strip():
                chunks.append({
                    "text": current_text.strip(),
                    "page": current_page,
                    "chunk_type": current_type,
                })
            current_text = line + "\n"
            current_type = "heading"
            continue

        # 偵測表格
        if stripped.startswith("|") and "|" in stripped[1:]:
            if current_type != "table":
                if current_text.strip():
                    chunks.append({
                        "text": current_text.strip(),
                        "page": current_page,
                        "chunk_type": current_type,
                    })
                current_text = ""
                current_type = "table"
            current_text += line + "\n"
            continue

        # 偵測程式碼區塊
        if stripped.startswith("```"):
            if current_type == "code":
                current_text += line + "\n"
                chunks.append({
                    "text": current_text.strip(),
                    "page": current_page,
                    "chunk_type": "code",
                })
                current_text = ""
                current_type = "text"
            else:
                if current_text.strip():
                    chunks.append({
                        "text": current_text.strip(),
                        "page": current_page,
                        "chunk_type": current_type,
                    })
                current_text = line + "\n"
                current_type = "code"
            continue

        # 偵測清單
        if stripped and (
            stripped[0] in "-*+"
            or (stripped[0].isdigit() and "." in stripped[:4])
        ):
            if current_type not in ("list", "code"):
                if current_text.strip() and current_type != "list":
                    chunks.append({
                        "text": current_text.strip(),
                        "page": current_page,
                        "chunk_type": current_type,
                    })
                    current_text = ""
                current_type = "list"
            current_text += line + "\n"
            continue

        # 偵測頁面分隔
        if stripped in ("---", "***", "___"):
            if current_text.strip():
                chunks.append({
                    "text": current_text.strip(),
                    "page": current_page,
                    "chunk_type": current_type,
                })
                current_text = ""
            current_page += 1
            current_type = "text"
            continue

        # 從非 table/list/code 回歸 text
        if current_type in ("table", "list") and not stripped:
            if current_text.strip():
                chunks.append({
                    "text": current_text.strip(),
                    "page": current_page,
                    "chunk_type": current_type,
                })
                current_text = ""
            current_type = "text"
            continue

        current_text += line + "\n"

    # 最後一段
    if current_text.strip():
        chunks.append({
            "text": current_text.strip(),
            "page": current_page,
            "chunk_type": current_type,
        })

    # 合併過小的 chunks（< 50 字）
    merged: list[dict[str, Any]] = []
    for chunk in chunks:
        if (
            merged
            and len(chunk["text"]) < 50
            and merged[-1]["chunk_type"] == chunk["chunk_type"]
        ):
            merged[-1]["text"] += "\n" + chunk["text"]
        else:
            merged.append(chunk)

    return merged
